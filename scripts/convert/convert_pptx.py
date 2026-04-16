#!/usr/bin/env python3
"""Convert PowerPoint .pptx files to Markdown.

Usage:
    python convert_pptx.py <input_file> <output_md> [--images-dir <dir>]

Exit code 0 on success, non-zero on failure.
Prints JSON summary to stdout on completion.
"""

import argparse
import io
import json
import os
import re
import sys
import uuid
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches
from PIL import Image


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def get_image_extension(image_data: bytes) -> str:
    """Determine image format from raw bytes."""
    try:
        img = Image.open(io.BytesIO(image_data))
        fmt = (img.format or "png").lower()
        return "jpg" if fmt == "jpeg" else fmt.strip()
    except Exception:
        return "png"


def validate_image(path: Path) -> bool:
    """Return True if *path* is a valid image file."""
    try:
        with Image.open(path) as img:
            img.verify()
        return True
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Slide title detection
# ---------------------------------------------------------------------------

def get_slide_title(slide) -> Optional[str]:
    """Try to extract a title from the slide's placeholder shapes.

    Returns the title text or None if no title placeholder is found.
    """
    try:
        if slide.shapes.title is not None:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
    except Exception:
        pass

    # Fallback: look for placeholder with idx 0 (title) or idx 1 (subtitle used as title)
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.idx in (0, 1):
                text = shape.text.strip()
                if text:
                    return text
        except Exception:
            continue

    return None


# ---------------------------------------------------------------------------
# Table conversion
# ---------------------------------------------------------------------------

def convert_table_to_markdown(table) -> List[str]:
    """Convert a pptx Table object to Markdown table lines."""
    md: List[str] = []
    try:
        rows = table.rows
        if not rows:
            return []

        # Build a 2-D list of cell texts
        table_data: List[List[str]] = []
        for row in rows:
            row_cells = [cell.text.strip().replace('\n', '<br>').replace('|', '\\|')
                         for cell in row.cells]
            table_data.append(row_cells)

        if not table_data:
            return ["*[Empty table]*"]

        # Normalise column count
        max_cols = max(len(r) for r in table_data)
        for i, row in enumerate(table_data):
            if len(row) < max_cols:
                table_data[i] = row + [""] * (max_cols - len(row))

        # Header row
        if len(table_data) > 1:
            headers = table_data[0]
            data_rows = table_data[1:]
        else:
            headers = [f"Col{i+1}" for i in range(max_cols)]
            data_rows = table_data

        md.append("| " + " | ".join(h or " " for h in headers) + " |")
        md.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in data_rows:
            md.append("| " + " | ".join(c or " " for c in row) + " |")

    except Exception:
        return ["*[Table parsing failed]*"]

    return md


# ---------------------------------------------------------------------------
# Shape text extraction
# ---------------------------------------------------------------------------

def extract_shape_text(shape) -> List[str]:
    """Extract Markdown-formatted text lines from a shape's text frame."""
    lines: List[str] = []
    if not shape.has_text_frame:
        return lines

    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)

    return lines


# ---------------------------------------------------------------------------
# Main conversion
# ---------------------------------------------------------------------------

def convert_pptx(input_file: str, output_md: str, images_dir: str) -> Dict:
    """Convert *input_file* (.pptx) → Markdown written to *output_md*.

    Returns a summary dict.
    """
    input_path = Path(input_file)
    output_path = Path(output_md)
    img_dir = Path(images_dir)

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    prs = Presentation(str(input_path))
    total_slides = len(prs.slides)

    # Ensure images directory exists
    img_dir.mkdir(parents=True, exist_ok=True)

    # Compute the relative path from the output markdown file to the images dir
    try:
        img_rel = os.path.relpath(img_dir, output_path.parent)
    except ValueError:
        # On Windows, relpath fails across drives
        img_rel = str(img_dir)

    markdown_content: List[str] = []
    images_info: List[Dict] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_lines: List[str] = []
        title_text = get_slide_title(slide)

        # Track which shapes have been used as the title so we don't duplicate
        title_shape_id = None
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title_shape_id = slide.shapes.title.shape_id
        except Exception:
            pass

        for shape in slide.shapes:
            # --- Images (shape_type 13 == PICTURE) ---
            if shape.shape_type == 13:
                try:
                    image = shape.image
                    image_bytes = image.blob

                    if len(image_bytes) < 100:
                        continue

                    ext = get_image_extension(image_bytes)
                    filename = f"{uuid.uuid4()}.{ext}"
                    dest = img_dir / filename

                    with open(dest, 'wb') as f:
                        f.write(image_bytes)

                    if validate_image(dest):
                        rel_img_path = os.path.join(img_rel, filename).replace('\\', '/')
                        slide_lines.append(f"![Image]({rel_img_path})")
                        images_info.append({
                            "filename": filename,
                            "path": str(dest),
                            "size": len(image_bytes),
                            "format": ext,
                        })
                    else:
                        dest.unlink(missing_ok=True)
                except Exception:
                    continue

            # --- Tables ---
            elif shape.has_table:
                try:
                    table_md = convert_table_to_markdown(shape.table)
                    if table_md:
                        slide_lines.append("")
                        slide_lines.extend(table_md)
                        slide_lines.append("")
                except Exception:
                    slide_lines.append("*[Table parsing failed]*")

            # --- Text shapes ---
            elif hasattr(shape, "text") and shape.text.strip():
                # Skip the title shape — we already use it as a heading
                if title_shape_id is not None and shape.shape_id == title_shape_id:
                    continue

                text_lines = extract_shape_text(shape)
                if text_lines:
                    slide_lines.extend(text_lines)

        # Skip slides with absolutely no content
        if not slide_lines and not title_text:
            continue

        # Slide heading
        if title_text:
            markdown_content.append(f"## Slide {slide_idx}: {title_text}")
        else:
            markdown_content.append(f"## Slide {slide_idx}")
        markdown_content.append("")

        if slide_lines:
            markdown_content.extend(slide_lines)
        else:
            markdown_content.append("*Slide has no text content*")
        markdown_content.append("")

    result_md = clean_markdown(markdown_content)

    if not result_md.strip():
        result_md = "*Presentation content is empty or cannot be parsed*"

    # Write output
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(result_md, encoding='utf-8')

    return {
        "success": True,
        "file_type": "pptx",
        "slides": total_slides,
        "images": len(images_info),
    }


# ---------------------------------------------------------------------------
# Markdown cleanup
# ---------------------------------------------------------------------------

def clean_markdown(lines: List[str]) -> str:
    """Join *lines* into a single Markdown string, trimming excess blanks."""
    content = "\n".join(lines)
    content = re.sub(r'\n{3,}', '\n\n', content)
    stripped = [l.rstrip() for l in content.split('\n')]
    while stripped and not stripped[0].strip():
        stripped.pop(0)
    while stripped and not stripped[-1].strip():
        stripped.pop()
    return '\n'.join(stripped)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(
        description="Convert a .pptx file to Markdown."
    )
    parser.add_argument("input_file", help="Path to the .pptx file")
    parser.add_argument("output_md", help="Path to write the resulting .md file")
    parser.add_argument(
        "--images-dir",
        default=None,
        help="Directory to save extracted images (default: same dir as output_md)",
    )
    args = parser.parse_args()

    images_dir = args.images_dir
    if images_dir is None:
        images_dir = str(Path(args.output_md).parent)

    try:
        summary = convert_pptx(args.input_file, args.output_md, images_dir)
        print(json.dumps(summary))
        return 0
    except Exception as exc:
        err = {"success": False, "error": str(exc)}
        print(json.dumps(err), file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
